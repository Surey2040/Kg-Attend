import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

# Define Theme Colors
COLOR_BLUE = RGBColor(74, 144, 226)
COLOR_GREEN = RGBColor(80, 200, 120)
COLOR_PURPLE = RGBColor(155, 89, 182)
COLOR_ORANGE = RGBColor(243, 156, 18)
COLOR_RED = RGBColor(231, 76, 60)
COLOR_DARKBLUE = RGBColor(44, 62, 80)
COLOR_GRAY = RGBColor(200, 200, 200)
COLOR_LIGHT_FILL = RGBColor(235, 240, 250)

# Load the V3 presentation to modify
prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v3.pptx')

def draw_light_box(slide, x, y, w, h, text, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_FILL
    shape.line.color.rgb = COLOR_DARKBLUE
    tf = shape.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = COLOR_DARKBLUE
    return shape

def draw_green_box(slide, x, y, w, h, text, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 245, 230)
    shape.line.color.rgb = RGBColor(39, 174, 96)
    tf = shape.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)
    tf.paragraphs[0].font.bold = True
    return shape

def draw_pipeline_arrow(slide, x_center, y_start, h=0.25, color=COLOR_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x_center - 0.075), Inches(y_start), Inches(0.15), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def draw_line_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# SLIDE 11 is at index 10
slide11 = prs.slides[10]

# Remove all shapes on slide 11 except the title
shapes_to_delete = []
for shape in slide11.shapes:
    if shape != slide11.shapes.title:
        shapes_to_delete.append(shape)

for shape in shapes_to_delete:
    shape.element.getparent().remove(shape.element)


left_texts = [
    "1. Authenticate Student JWT Token",
    "2. Verify Active Account & Batch Registration",
    "3. Validate Attendance Session Existence",
    "4. Confirm Session Status is 'ACTIVE'",
    "5. Extract Cryptographic QR Nonce Payload",
    "6. Validate HMAC-SHA256 Signature Integrity",
    "7. Verify 10-Second Time Window Freshness"
]

right_texts = [
    "8. Check Redis Nonce Reuse (Anti-Replay)",
    "9. Verify Student Enrolled in Course/Batch",
    "10. Extract Device GPS Coordinates & Accuracy",
    "11. Compute Haversine Distance vs Room Lat/Lng",
    "12. Validate Campus Geofence Boundary (≤25m)",
]

box_h = 0.5
box_w = 5.0
gap = 0.25
left_x = 1.0
right_x = 7.0

# Draw Left Column (Boxes 1 to 7)
for i, text in enumerate(left_texts):
    y = 1.6 + i * (box_h + gap)
    draw_light_box(slide11, left_x, y, box_w, box_h, text, font_size=11)
    # Draw arrow down to next box
    if i < 6:
        draw_pipeline_arrow(slide11, left_x + box_w/2, y + box_h, h=gap, color=COLOR_BLUE)

# Draw Right Column (Boxes 8 to 12)
for i, text in enumerate(right_texts):
    y = 1.6 + i * (box_h + gap)
    draw_light_box(slide11, right_x, y, box_w, box_h, text, font_size=11)
    # Draw arrow down
    draw_pipeline_arrow(slide11, right_x + box_w/2, y + box_h, h=gap, color=COLOR_BLUE)

# Box 13 (Green)
y_13 = 1.6 + 5 * (box_h + gap)
draw_green_box(slide11, right_x, y_13, box_w, box_h, "13. Atomically Commit AttendanceRecord to DB", font_size=11)

# Connection from 7 to 8
bottom_y = 1.6 + 6 * (box_h + gap) + box_h # 6.6
center_left = left_x + box_w/2
center_right = right_x + box_w/2

draw_line_rect(slide11, center_left - 0.025, bottom_y, 0.05, 0.2, COLOR_BLUE) # down
draw_line_rect(slide11, center_left - 0.025, bottom_y + 0.2, center_right - center_left + 0.05, 0.05, COLOR_BLUE) # right
draw_line_rect(slide11, center_right - 0.025, 1.2, 0.05, bottom_y + 0.2 - 1.2, COLOR_BLUE) # up
draw_pipeline_arrow(slide11, center_right, 1.4, h=0.2, color=COLOR_BLUE) # arrow down to 8

# "continued" text
tb = slide11.shapes.add_textbox(Inches(5.8), Inches(6.85), Inches(1.5), Inches(0.2)).text_frame
tb.text = "continued →"
tb.paragraphs[0].font.size = Pt(9)
tb.paragraphs[0].font.italic = True
tb.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
tb.paragraphs[0].alignment = PP_ALIGN.CENTER

# Green arrow out of 13
draw_pipeline_arrow(slide11, center_right, y_13 + box_h, h=gap, color=RGBColor(39, 174, 96))

# SUCCESS text
tb_success = slide11.shapes.add_textbox(Inches(right_x), Inches(y_13 + box_h + gap), Inches(box_w), Inches(0.5)).text_frame
tb_success.text = "SUCCESS → Real-Time Socket.IO Telemetry\nPushed to Faculty Dashboard"
tb_success.paragraphs[0].font.bold = True
tb_success.paragraphs[0].font.size = Pt(11)
tb_success.paragraphs[0].font.color.rgb = RGBColor(39, 174, 96)
tb_success.paragraphs[0].alignment = PP_ALIGN.CENTER
tb_success.paragraphs[1].font.bold = True
tb_success.paragraphs[1].font.size = Pt(11)
tb_success.paragraphs[1].font.color.rgb = RGBColor(39, 174, 96)
tb_success.paragraphs[1].alignment = PP_ALIGN.CENTER

# Bottom Figure caption
tb_fig = slide11.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.4)).text_frame
tb_fig.text = "Fig. 5.1 — 13-Step zero-trust validation pipeline executed per scan request"
tb_fig.paragraphs[0].font.italic = True
tb_fig.paragraphs[0].font.size = Pt(13)
tb_fig.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
tb_fig.paragraphs[0].alignment = PP_ALIGN.CENTER

prs.save(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v4.pptx')
print("Successfully modified PresenceIQ_MCA_Review1_Final_v4.pptx (15 slides) - Recreated 13-stage visual pipeline.")
