import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

# Define Theme Colors
COLOR_BLUE = RGBColor(74, 144, 226)
COLOR_GREEN = RGBColor(80, 200, 120)
COLOR_PURPLE = RGBColor(155, 89, 182)
COLOR_ORANGE = RGBColor(243, 156, 18)
COLOR_RED = RGBColor(231, 76, 60)
COLOR_DARKBLUE = RGBColor(44, 62, 80)
COLOR_GRAY = RGBColor(200, 200, 200)
COLOR_LIGHT_FILL = RGBColor(235, 240, 250)

# Load the V5 presentation
prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v5.pptx')
layout_title_content = prs.slide_layouts[1]

def draw_box(slide, x, y, w, h, text, color=COLOR_BLUE, font_size=12, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_DARKBLUE
    tf = shape.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.bold = bold
    return shape

def draw_arrow_down(slide, x, y, w, h, color=COLOR_DARKBLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def draw_arrow_right(slide, x, y, w, h, color=COLOR_DARKBLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def add_simple_slide(title_text, body_text, notes=""):
    slide = prs.slides.add_slide(layout_title_content)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = body_text
    if slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes_text
    return slide

def replace_text_in_shape(shape, search_text, replace_text):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if search_text in paragraph.text:
            paragraph.runs[0].text = paragraph.text.replace(search_text, replace_text)
            for i in range(len(paragraph.runs)-1, 0, -1):
                p = paragraph._p
                p.remove(paragraph.runs[i]._r)

# 1. Update Contents List on Slide 2 (index 1)
slide2 = prs.slides[1]
contents_list = (
    "1. Introduction and Problem Definition\n2. Existing System and Limitations\n3. Proposed PresenceIQ System\n"
    "4. Overall Technology Architecture\n5. Dynamic QR and HMAC Security\n6. Redis Replay Protection\n"
    "7. Device Binding and Proxy Prevention\n8. GPS and Geofence Verification\n9. 13-Stage Validation Pipeline\n"
    "10. YOLO ML Headcount Verification\n11. Real-Time and API Security\n12. Docker Deployment Architecture\n13. System Requirements\n14. Conclusion"
)
for shape in slide2.shapes:
    if shape.has_text_frame and "INTRODUCTION AND PROBLEM" in shape.text.upper():
        shape.text_frame.text = contents_list

# 2. Add YOLO ML Slide (Will be at the end initially)
slide_yolo = add_simple_slide("YOLO ML HEADCOUNT VERIFICATION (AI PROXY DETECTION)", "")
sp = slide_yolo.placeholders[1].element
sp.getparent().remove(sp)

# Row 1
draw_box(slide_yolo, 1.0, 1.8, 3.5, 1.2, "1. Faculty Mobile App\n\nCaptures wide-angle classroom photo during active session", COLOR_PURPLE, 12)
draw_arrow_right(slide_yolo, 4.7, 2.2, 0.8, 0.4, COLOR_DARKBLUE)

draw_box(slide_yolo, 5.7, 1.8, 4.0, 1.2, "2. YOLOv8 ML Model Engine\n\nProcesses image and accurately counts physical student headcount", COLOR_BLUE, 12)
draw_arrow_down(slide_yolo, 7.5, 3.2, 0.4, 0.6, COLOR_DARKBLUE)

# Row 2
draw_box(slide_yolo, 5.7, 4.0, 4.0, 1.2, "3. Validation Server\n\nCompares Physical Headcount vs Marked QR Attendances", COLOR_PURPLE, 12)

# Arrows from Validation server to the decision boxes
draw_arrow_down(slide_yolo, 6.0, 5.4, 0.3, 0.5, COLOR_GREEN)
draw_arrow_down(slide_yolo, 9.1, 5.4, 0.3, 0.5, COLOR_RED)

# Row 3 (Decisions)
draw_box(slide_yolo, 4.0, 6.0, 4.2, 1.0, "MATCH / NORMAL\nPhysical Heads >= QR Scans\nAll attendances verified authentic.", COLOR_GREEN, 11)
draw_box(slide_yolo, 8.5, 6.0, 4.2, 1.0, "MISMATCH / ALERT\nPhysical Heads < QR Scans\nProxy or Geofence Spoofing Detected!", COLOR_RED, 11)

tb_note = slide_yolo.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4.0), Inches(1.5)).text_frame
tb_note.text = "How it catches proxies:\nIf a student marks attendance from the canteen using a forwarded QR code, the QR count becomes 50, but YOLO only sees 49 heads in class. The system instantly red-flags the discrepancy."
tb_note.paragraphs[0].font.bold = True
tb_note.paragraphs[0].font.size = Pt(13)
tb_note.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
for p in tb_note.paragraphs[1:]:
    p.font.size = Pt(11)

# 3. Move the newly appended slide (index 15) to index 11 (Slide 12)
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
# The appended slide is at the very end
yolo_slide_xml = slides[-1]
xml_slides.remove(yolo_slide_xml)
# Insert at index 11
xml_slides.insert(11, yolo_slide_xml)

prs.save(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v6.pptx')
print("Successfully generated PresenceIQ_MCA_Review1_Final_v6.pptx (16 slides) - Appended and moved YOLO slide to Slide 12.")
