from pptx import Presentation

prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\MCA - Review1  PPT Template.pptx')

def replace_text_in_shape(shape, search_text, replace_text):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if search_text in paragraph.text:
            paragraph.runs[0].text = paragraph.text.replace(search_text, replace_text)
            for i in range(len(paragraph.runs)-1, 0, -1):
                p = paragraph._p
                p.remove(paragraph.runs[i]._r)

slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, "Your Project Full Title", "PresenceIQ: Real-Time Geofenced & Dynamic QR Smart Attendance Engine")
        replace_text_in_shape(shape, "Your Name  Reg.No", "Surender vignesh M    Reg.No: 2538M0163")
        replace_text_in_shape(shape, "Your Guide Name with Qualification  Designation", "Dr.G.Rajesh")

slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, "Your Project Full Title", "PresenceIQ")

def add_slide(title_text, body_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = title_text
    body.text = body_text

def add_two_column_slide(title_text, left_text, right_text):
    slide_layout = prs.slide_layouts[3] # Two Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    left_body = slide.placeholders[1]
    right_body = slide.placeholders[2]
    
    title.text = title_text
    left_body.text = left_text
    right_body.text = right_text

intro_text = (
    "Traditional attendance systems (roll-calls, plain QR codes) are slow and proxy-prone.\n"
    "Hardware-based biometrics are expensive and create bottlenecks.\n"
    "Need for a zero-hardware, real-time solution for universities and enterprise workspaces.\n"
    "PresenceIQ aims to provide a cryptographically secure, sub-second latency attendance engine."
)
add_slide("INTRODUCTION (Problem Definition)", intro_text)

existing_text = (
    "Static QR Codes: Easily shared via WhatsApp, leading to mass proxies.\n"
    "RFID / Biometrics: Hardware-dependent, high maintenance cost, and slow processing.\n"
    "Lack of Geofencing: Students can mark attendance from outside the classroom.\n"
    "No Real-Time Monitoring: Delays in updating attendance records and identifying defaulters."
)
add_slide("EXISTING SYSTEM (Limitations)", existing_text)

proposed_text = (
    "Objective: Real-Time Geofenced & Dynamic QR Smart Attendance Engine.\n"
    "Methodology: 13-Step Zero-Trust Security Pipeline ensuring zero proxies.\n"
    "Dynamic QR generation with 10-second TTL and HMAC-SHA256 signatures.\n"
    "GPS Geofencing (Haversine formula within 25m) & Device Hardware ID Binding.\n"
    "3-Container Architecture (Node.js backend, Redis for replay-locks, Postgres vault)."
)
add_slide("PROPOSED SYSTEM", proposed_text)

# Add PIPELINE SLIDE
pipeline_left = (
    "1. Authenticate Student JWT Token\n"
    "2. Verify Active Account & Batch Registration\n"
    "3. Validate Attendance Session Existence\n"
    "4. Confirm Session Status is 'ACTIVE'\n"
    "5. Extract Cryptographic QR Nonce Payload\n"
    "6. Validate HMAC-SHA256 Signature Integrity\n"
    "7. Verify 10-Second Time Window Freshness"
)

pipeline_right = (
    "8. Check Redis Nonce Reuse (Anti-Replay)\n"
    "9. Verify Student Enrolled in Course/Batch\n"
    "10. Extract Device GPS Coordinates & Accuracy\n"
    "11. Compute Haversine Distance vs Room Lat/Lng\n"
    "12. Validate Campus Geofence Boundary (≤25m)\n"
    "13. Atomically Commit AttendanceRecord to DB"
)

add_two_column_slide("13-STEP SECURITY VALIDATION PIPELINE", pipeline_left, pipeline_right)

sys_req_text = (
    "Hardware Specifications:\n"
    "Client: Smartphone with Camera and GPS.\n"
    "Server: Standard Cloud Instance or Local Server.\n\n"
    "Software Specifications:\n"
    "Frontend: React.js, Vite, Tailwind CSS, Socket.IO Client.\n"
    "Backend: Node.js, Express.js, Socket.IO, Prisma ORM.\n"
    "Database & Cache: PostgreSQL, Redis 7.\n"
    "Deployment: Docker, Docker-Compose."
)
add_slide("SYSTEM REQUIREMENT", sys_req_text)

conclusion_text = (
    "PresenceIQ completely eliminates proxy attendance using dynamic cryptographic QRs and strict GPS geofencing.\n"
    "High scalability and low latency achieved through a modern 3-container microservice architecture.\n"
    "Ready to be deployed as a zero-hardware, cost-effective solution for modern educational institutions."
)
add_slide("CONCLUSIONS", conclusion_text)

prs.save(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_Review1_v3.pptx')
print("Saved PresenceIQ_Review1_v3.pptx")
