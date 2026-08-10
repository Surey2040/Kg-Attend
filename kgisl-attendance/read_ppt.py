from pptx import Presentation

prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\MCA - Review1  PPT Template.pptx')

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(f"Text: {shape.text.replace(chr(10), ' ')}")
