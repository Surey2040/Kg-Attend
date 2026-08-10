import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

# Define Theme Colors
COLOR_PURPLE = RGBColor(155, 89, 182)
COLOR_DARKBLUE = RGBColor(44, 62, 80)
COLOR_GREEN = RGBColor(39, 174, 96)

# Load the V6 presentation
prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v6.pptx')

# 1. Add JSON example to Slide 7 (index 6) - HMAC
slide7 = prs.slides[6]
json_hmac_shape = slide7.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
json_hmac = json_hmac_shape.text_frame
json_hmac.text = 'JSON Payload Example: { "sessionId": "c8f2b1", "nonce": "xyz987", "hmacSignature": "d41d8cd98f00b204e9800998ecf8427e" }'
json_hmac.paragraphs[0].font.name = "Consolas"
json_hmac.paragraphs[0].font.size = Pt(14)
json_hmac.paragraphs[0].font.color.rgb = COLOR_GREEN

# 2. Add JSON example to Slide 10 (index 9) - GPS
slide10 = prs.slides[9]

# Narrow the existing text placeholder to prevent overlap
for shape in slide10.shapes:
    if shape.is_placeholder and shape.placeholder_format.idx == 1:
        shape.width = Inches(7.0)

# Add the JSON Box on the right side
json_gps_shape = slide10.shapes.add_textbox(Inches(7.5), Inches(2.0), Inches(5.0), Inches(4.0))
json_gps = json_gps_shape.text_frame

json_gps.text = "API Request JSON Payload:\n\n" \
          "{\n" \
          "  \"studentId\": \"2538M0163\",\n" \
          "  \"location\": {\n" \
          "    \"latitude\": 11.0772,\n" \
          "    \"longitude\": 76.9902,\n" \
          "    \"accuracy\": 15.2\n" \
          "  },\n" \
          "  \"qrPayload\": {\n" \
          "    \"sessionId\": \"sess_8x9\",\n" \
          "    \"nonce\": \"xyz123\"\n" \
          "  }\n" \
          "}"

json_gps.paragraphs[0].font.bold = True
json_gps.paragraphs[0].font.size = Pt(16)
json_gps.paragraphs[0].font.color.rgb = COLOR_DARKBLUE

for i in range(2, len(json_gps.paragraphs)):
    json_gps.paragraphs[i].font.name = "Consolas"
    json_gps.paragraphs[i].font.size = Pt(15)
    json_gps.paragraphs[i].font.color.rgb = COLOR_PURPLE

# Add a neat background to the JSON box to make it pop
json_gps_shape.fill.solid()
json_gps_shape.fill.fore_color.rgb = RGBColor(245, 245, 250)
json_gps_shape.line.color.rgb = COLOR_PURPLE

prs.save(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v7.pptx')
print("Successfully generated PresenceIQ_MCA_Review1_Final_v7.pptx - Added JSON payload examples to Slide 7 and Slide 10.")
