import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

# Define Theme Colors
COLOR_BLUE = RGBColor(74, 144, 226)
COLOR_GREEN = RGBColor(80, 200, 120)
COLOR_PURPLE = RGBColor(155, 89, 182)
COLOR_ORANGE = RGBColor(243, 156, 18)
COLOR_RED = RGBColor(231, 76, 60)
COLOR_DARKBLUE = RGBColor(44, 62, 80)
COLOR_GRAY = RGBColor(200, 200, 200)

prs = Presentation(r'c:\Users\surreyy\Desktop\presenceIQ\MCA - Review1  PPT Template.pptx')
layout_title_content = prs.slide_layouts[1]
layout_two_content = prs.slide_layouts[3]
layout_blank = prs.slide_layouts[6]

def replace_text_in_shape(shape, search_text, replace_text):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if search_text in paragraph.text:
            paragraph.runs[0].text = paragraph.text.replace(search_text, replace_text)
            for i in range(len(paragraph.runs)-1, 0, -1):
                p = paragraph._p
                p.remove(paragraph.runs[i]._r)

def add_speaker_notes(slide, notes_text):
    if slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes_text

def draw_box(slide, x, y, w, h, text, color=COLOR_BLUE, font_size=12, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_DARKBLUE
    tf = shape.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.paragraphs[0].font.bold = bold
    return shape

def draw_arrow_down(slide, x, y, w, h, color=COLOR_DARKBLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def draw_arrow_right(slide, x, y, w, h, color=COLOR_DARKBLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def add_simple_slide(title_text, body_text, notes=""):
    slide = prs.slides.add_slide(layout_title_content)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = body_text
    add_speaker_notes(slide, notes)
    return slide

def add_two_column_slide(title_text, left_text, right_text, notes=""):
    slide = prs.slides.add_slide(layout_two_content)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = left_text
    slide.placeholders[2].text = right_text
    add_speaker_notes(slide, notes)
    return slide

# SLIDE 1
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, "Your Project Full Title", "PresenceIQ: Real-Time Geofenced & Dynamic QR Smart Attendance Engine\nMCA Project Review – I")
        replace_text_in_shape(shape, "Your Name  Reg.No", "Surender vignesh M    Reg.No: 2538M0163")
        replace_text_in_shape(shape, "Your Guide Name with Qualification  Designation", "Dr.G.Rajesh")

# SLIDE 2 - Updated Contents List (Removed unused sections)
slide2 = prs.slides[1]
contents_list = (
    "1. Introduction and Problem Definition\n2. Existing System and Limitations\n3. Proposed PresenceIQ System\n"
    "4. Overall Technology Architecture\n5. Dynamic QR and HMAC Security\n6. Redis Replay Protection\n"
    "7. Device Binding and Proxy Prevention\n8. GPS and Geofence Verification\n9. 13-Stage Validation Pipeline\n"
    "10. Real-Time and API Security\n11. Docker Deployment Architecture\n12. System Requirements\n13. Conclusion"
)
for shape in slide2.shapes:
    if shape.has_text_frame:
        replace_text_in_shape(shape, "Your Project Full Title", "PresenceIQ")
        if "INTRODUCTION" in shape.text:
            shape.text_frame.text = contents_list

# SLIDE 3
add_simple_slide("INTRODUCTION AND PROBLEM DEFINITION", 
    "• Manual attendance consumes excessive classroom time.\n"
    "• Static QR codes can be easily photographed and forwarded.\n"
    "• Credentials are often shared for proxy attendance.\n"
    "• Basic QR systems do not prove physical presence.\n"
    "• Delayed reporting prevents real-time faculty intervention.\n"
    "• Paper records are difficult to audit securely.\n\n"
    "Core philosophy: PresenceIQ treats every scan as an untrusted security request that must mathematically prove identity, QR authenticity, freshness, device continuity, and physical presence."
)

# SLIDE 4
add_two_column_slide("EXISTING SYSTEM AND LIMITATIONS",
    "Existing Systems:\n• Manual roll call\n• Biometric devices\n• RFID cards\n• Static QR attendance\n• GPS-only attendance",
    "Limitations:\n• Time-consuming\n• Dedicated hardware cost\n• Card or credential sharing\n• QR screenshot forwarding\n• GPS spoofing\n• Poor auditability\n• No real-time dashboard"
)

# SLIDE 5 (Overview)
slide5 = prs.slides.add_slide(layout_blank)
tx = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
tx.text_frame.text = "PROPOSED SYSTEM OVERVIEW"
tx.text_frame.paragraphs[0].font.size = Pt(32)
draw_box(slide5, Inches(1), Inches(2.0), Inches(8), Inches(0.8), "Client Layer:\nStudent React PWA | Faculty Dashboard | Admin Dashboard", COLOR_BLUE, 14)
draw_arrow_down(slide5, Inches(4.8), Inches(2.8), Inches(0.4), Inches(0.4), COLOR_DARKBLUE)
draw_box(slide5, Inches(1), Inches(3.2), Inches(8), Inches(0.8), "API and Real-Time Layer:\nExpress REST API | JWT Authentication | Socket.IO", COLOR_PURPLE, 14)
draw_arrow_down(slide5, Inches(4.8), Inches(4.0), Inches(0.4), Inches(0.4), COLOR_DARKBLUE)
draw_box(slide5, Inches(1), Inches(4.4), Inches(8), Inches(0.8), "Security and Business Layer:\nQR Service | Validation Pipeline | Device Binding | Geofence Engine", COLOR_GREEN, 14)
draw_arrow_down(slide5, Inches(4.8), Inches(5.2), Inches(0.4), Inches(0.4), COLOR_DARKBLUE)
draw_box(slide5, Inches(1), Inches(5.6), Inches(8), Inches(0.8), "Data Layer:\nRedis | PostgreSQL through Prisma", COLOR_DARKBLUE, 14)

# SLIDE 6 (Tech Stack)
slide6 = add_simple_slide("TECHNOLOGY STACK: WHAT, WHY AND HOW", "")
sp = slide6.placeholders[1].element
sp.getparent().remove(sp)
tbl = slide6.shapes.add_table(7, 4, Inches(0.5), Inches(1.5), Inches(12.0), Inches(5.0)).table
headers = ["Technology", "What", "Why Selected", "How Used"]
for i, h in enumerate(headers): tbl.cell(0, i).text = h
data = [
    ["React 18 & Vite", "Frontend FW", "Fast HMR, modern state", "Student PWA & Admin UI"],
    ["Node.js & Express", "API Backend", "Async I/O, scalable", "Core Validation Engine"],
    ["Socket.IO", "WebSocket", "Real-time bi-directional", "Live faculty telemetry"],
    ["PostgreSQL 16", "RDBMS", "ACID compliance", "Durable attendance vault"],
    ["Redis 7", "In-Memory Store", "Sub-millisecond speed", "Replay locks, rate limits"],
    ["Docker Compose", "Orchestration", "Isolated reproducibility", "Multi-container server"]
]
for r, row in enumerate(data):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
        tbl.cell(r+1, c).text_frame.paragraphs[0].font.size = Pt(14)

# SLIDE 7 (HMAC QR) - REDESIGNED TO PREVENT OVERLAP AND SHOW ARROWS
slide7 = add_simple_slide("DYNAMIC QR AND HMAC-SHA256 SECURITY", "")
slide7.placeholders[1].element.getparent().remove(slide7.placeholders[1].element)

draw_box(slide7, Inches(0.5), Inches(2.0), Inches(2.5), Inches(0.4), "Faculty Dashboard", COLOR_GRAY, 12)
draw_box(slide7, Inches(3.5), Inches(2.0), Inches(3.0), Inches(0.4), "Node.js Backend", COLOR_PURPLE, 12)
draw_box(slide7, Inches(7.0), Inches(2.0), Inches(2.5), Inches(0.4), "Redis/PostgreSQL", COLOR_DARKBLUE, 12)

# Start session
draw_box(slide7, Inches(0.5), Inches(2.8), Inches(2.5), Inches(0.6), "Start Attendance Session", COLOR_BLUE)
draw_arrow_right(slide7, Inches(3.0), Inches(2.9), Inches(0.5), Inches(0.4), COLOR_DARKBLUE)
# Backend generate
draw_box(slide7, Inches(3.5), Inches(2.8), Inches(3.0), Inches(0.6), "Generate 256-bit token & Nonce", COLOR_BLUE)
draw_arrow_down(slide7, Inches(4.8), Inches(3.4), Inches(0.4), Inches(0.4), COLOR_DARKBLUE)
# Backend HMAC
draw_box(slide7, Inches(3.5), Inches(3.8), Inches(3.0), Inches(0.6), "HMAC-SHA256(Secret, Payload)", COLOR_BLUE)
# To Redis
draw_arrow_right(slide7, Inches(6.5), Inches(3.9), Inches(0.5), Inches(0.4), COLOR_DARKBLUE)
draw_box(slide7, Inches(7.0), Inches(3.8), Inches(2.5), Inches(0.6), "Store TokenHash & Nonce", COLOR_BLUE)
# To Faculty
draw_arrow_down(slide7, Inches(4.8), Inches(4.4), Inches(0.4), Inches(0.4), COLOR_DARKBLUE)
draw_box(slide7, Inches(3.5), Inches(4.8), Inches(3.0), Inches(0.6), "Broadcast QR via Socket", COLOR_BLUE)
# faculty arrow down
draw_arrow_down(slide7, Inches(1.5), Inches(3.4), Inches(0.4), Inches(1.4), COLOR_DARKBLUE) 
draw_box(slide7, Inches(0.5), Inches(4.8), Inches(2.5), Inches(0.6), "Display Dynamic QR", COLOR_BLUE)

fb = slide7.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12), Inches(1)).text_frame
fb.text = "Signature = HMAC-SHA256(SecretKey, SessionID || Token || IssuedAt || ExpiresAt || Nonce)"
fb.paragraphs[0].font.bold = True
fb.paragraphs[0].font.size = Pt(20)

# SLIDE 8 (Redis)
slide8 = add_simple_slide("HOW REDIS PREVENTS QR REPLAY", "")
sp = slide8.placeholders[1].element
sp.getparent().remove(sp)
draw_box(slide8, Inches(1), Inches(2.5), Inches(2), Inches(0.8), "Request A", COLOR_GRAY)
draw_box(slide8, Inches(1), Inches(4.5), Inches(2), Inches(0.8), "Request B", COLOR_GRAY)

draw_arrow_right(slide8, Inches(3.0), Inches(2.7), Inches(1.0), Inches(0.4), COLOR_DARKBLUE)
draw_arrow_right(slide8, Inches(3.0), Inches(4.7), Inches(1.0), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide8, Inches(4), Inches(3.2), Inches(2.5), Inches(1.2), "Redis Lua Atomic Claim", COLOR_RED)

draw_arrow_right(slide8, Inches(6.5), Inches(2.7), Inches(1.0), Inches(0.4), COLOR_DARKBLUE)
draw_arrow_right(slide8, Inches(6.5), Inches(4.7), Inches(1.0), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide8, Inches(7.5), Inches(2.5), Inches(2), Inches(0.8), "SUCCESS", COLOR_GREEN)
draw_box(slide8, Inches(7.5), Inches(4.5), Inches(2), Inches(0.8), "REJECTED (Replay)", COLOR_RED)

code = slide8.shapes.add_textbox(Inches(2), Inches(5.5), Inches(8), Inches(1.5)).text_frame
code.text = "IF active QR exists:\n    SET replayKey studentId NX PX remainingTTL\n    success -> first claim\n    failure -> duplicate/replay"
code.paragraphs[0].font.size = Pt(16)
code.paragraphs[1].font.size = Pt(16)
code.paragraphs[2].font.size = Pt(16)
code.paragraphs[3].font.size = Pt(16)

# SLIDE 9 (Device Binding) - REDESIGNED AS IMAGE DIAGRAM
slide9 = add_simple_slide("DEVICE BINDING AND PROXY PREVENTION", "")
sp = slide9.placeholders[1].element
sp.getparent().remove(sp)

# TOP ROW: Registration
t1 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.5)).text_frame
t1.text = "First Login (Device Registration)"
t1.paragraphs[0].font.size = Pt(18)
t1.paragraphs[0].font.bold = True
t1.paragraphs[0].font.color.rgb = COLOR_PURPLE

draw_box(slide9, Inches(0.5), Inches(2.0), Inches(2.2), Inches(0.8), "Student Browser", COLOR_GRAY)
draw_arrow_right(slide9, Inches(2.7), Inches(2.2), Inches(0.6), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide9, Inches(3.3), Inches(2.0), Inches(2.2), Inches(0.8), "Generate random UUID\nSave in localStorage", COLOR_BLUE)
draw_arrow_right(slide9, Inches(5.5), Inches(2.2), Inches(0.6), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide9, Inches(6.1), Inches(2.0), Inches(2.2), Inches(0.8), "Send UUID with\nLogin Request", COLOR_BLUE)
draw_arrow_right(slide9, Inches(8.3), Inches(2.2), Inches(0.6), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide9, Inches(8.9), Inches(2.0), Inches(2.5), Inches(0.8), "Backend saves UUID to\nStudent.deviceId in DB", COLOR_PURPLE)

# BOTTOM ROW: Verification
t2 = slide9.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(4), Inches(0.5)).text_frame
t2.text = "Subsequent QR Scan (Verification)"
t2.paragraphs[0].font.size = Pt(18)
t2.paragraphs[0].font.bold = True
t2.paragraphs[0].font.color.rgb = COLOR_PURPLE

draw_box(slide9, Inches(0.5), Inches(4.5), Inches(2.2), Inches(0.8), "Read UUID from\nlocalStorage", COLOR_GRAY)
draw_arrow_right(slide9, Inches(2.7), Inches(4.7), Inches(0.6), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide9, Inches(3.3), Inches(4.5), Inches(2.2), Inches(0.8), "Send UUID with\nAttendance Request", COLOR_BLUE)
draw_arrow_right(slide9, Inches(5.5), Inches(4.7), Inches(0.6), Inches(0.4), COLOR_DARKBLUE)

draw_box(slide9, Inches(6.1), Inches(4.5), Inches(2.2), Inches(0.8), "Backend compares with\nDB Student.deviceId", COLOR_PURPLE)

# Branching
draw_arrow_right(slide9, Inches(8.3), Inches(4.2), Inches(0.6), Inches(0.4), COLOR_DARKBLUE) # Up-right
draw_arrow_right(slide9, Inches(8.3), Inches(5.2), Inches(0.6), Inches(0.4), COLOR_DARKBLUE) # Down-right

draw_box(slide9, Inches(8.9), Inches(3.9), Inches(2.5), Inches(0.6), "MATCH: Access Granted", COLOR_GREEN)
draw_box(slide9, Inches(8.9), Inches(4.9), Inches(2.5), Inches(0.6), "MISMATCH: Rejected", COLOR_RED)


# SLIDE 10 (GPS)
slide10 = add_simple_slide("GPS AND GEOFENCE VERIFICATION", 
    "Verification Flow:\n"
    "1. Student GPS Coordinates & Accuracy received.\n"
    "2. Validate reported accuracy against threshold.\n"
    "3. Calculate Haversine distance vs Room Coordinates.\n"
    "4. Compare with room geofence radius.\n"
    "5. Accept or reject request.\n\n"
    "Haversine Formula:\n"
    "d = 2 * r * arcsin( sqrt( sin²(Δφ/2) + cos(φ1)*cos(φ2)*sin²(Δλ/2) ) )\n\n"
    "Note: Room coordinates and radius are stored in PostgreSQL. Browser GPS can be spoofed, making it one layer in a defence-in-depth model."
)

# SLIDE 11 (13-Stage Pipeline)
slide11 = add_two_column_slide("13-STAGE ZERO-TRUST VALIDATION PIPELINE",
    "1. JWT Authentication\n2. Account and Role Validation\n3. Request Schema Validation (Zod)\n4. Session Existence Check\n5. Active Session Check\n6. HMAC Signature Verification\n7. QR Freshness Validation",
    "8. Redis Active-Token Validation\n9. Session Payload Matching\n10. Device Binding Check\n11. GPS and Accuracy Validation\n12. Haversine Geofence Check (≤25m)\n13. Atomic Record, Audit and Broadcast"
)
tb = slide11.shapes.add_textbox(Inches(1), Inches(6.5), Inches(10), Inches(0.5)).text_frame
tb.text = "Attendance is accepted only when all required trust gates succeed."
tb.paragraphs[0].font.bold = True

# SLIDE 12 (JWT)
slide12 = add_simple_slide("REAL-TIME AND API SECURITY",
    "API Security:\n"
    "• Short-lived access JWT & longer-lived refresh JWT.\n"
    "• Single-use refresh-token rotation & family revocation on reuse.\n"
    "• Role-based access control.\n"
    "• Authentication and scan rate limiting.\n"
    "• Zod request validation.\n"
    "• Passwords are hashed using bcrypt, not reversibly encrypted.\n\n"
    "Refresh Flow:\n"
    "Login -> Access + Refresh Token -> Access Expires -> Refresh Consumed -> New Pair Issued.\n"
    "Attack: Old Refresh Reused -> Revoke Family -> Force Re-login + Audit."
)

# SLIDE 13 (Docker)
slide13 = add_simple_slide("DOCKER MULTI-CONTAINER ARCHITECTURE", "")
sp = slide13.placeholders[1].element
sp.getparent().remove(sp)
draw_box(slide13, Inches(0.5), Inches(2.0), Inches(8.5), Inches(5), "", RGBColor(240, 240, 240))
tt = slide13.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(4), Inches(0.5)).text_frame
tt.text = "Single Cloud/College Server"

draw_box(slide13, Inches(1.0), Inches(2.7), Inches(7.5), Inches(1.0), "Container 1: Backend\nNode.js + Express + Socket.IO + Prisma", COLOR_BLUE)
draw_box(slide13, Inches(1.0), Inches(4.2), Inches(3.5), Inches(1.2), "Container 2: Redis 7\nQR TTL | Replay Locks", COLOR_RED)
draw_box(slide13, Inches(5.0), Inches(4.2), Inches(3.5), Inches(1.2), "Container 3: PostgreSQL 16\nUsers | Sessions | Audit Logs", COLOR_DARKBLUE)

draw_box(slide13, Inches(10.0), Inches(3.0), Inches(2.5), Inches(1.0), "React Frontend\n(Internet/Browser)", COLOR_GRAY)

# SLIDE 14 (Requirements)
slide14 = add_two_column_slide("SYSTEM REQUIREMENTS",
    "Software:\n• Node.js 20+\n• React 18, TypeScript\n• Express 4\n• PostgreSQL 16\n• Redis 7\n• Prisma ORM\n• Docker & Docker Compose\n• Modern browser (Crypto API)",
    "Hardware (Suggested sizing):\n\nDevelopment:\n• 4-core CPU, 8 GB RAM\n• 10 GB free storage\n\nPilot Deployment:\n• 4 vCPU, 8-16 GB RAM\n• SSD storage\n• HTTPS domain\n• Automated PG backups"
)

# SLIDE 15 (Conclusion)
slide15 = add_simple_slide("CONCLUSION",
    "PresenceIQ converts attendance from a simple QR scan into a layered zero-trust verification transaction.\n\n"
    "Zero dedicated hardware. Multiple independent trust signals. Real-time institutional visibility.\n\n"
)
eqn = slide15.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2)).text_frame
eqn.text = "Trusted Attendance = \nAuthenticated Identity + Authentic Fresh QR + Active Redis State + \nReplay Protection + Registered Device + Verified Location + \nAtomic Database Record + Real-Time Audit Evidence"
eqn.paragraphs[0].font.bold = True
eqn.paragraphs[0].font.size = Pt(24)
eqn.paragraphs[0].font.color.rgb = COLOR_DARKBLUE

prs.save(r'c:\Users\surreyy\Desktop\presenceIQ\PresenceIQ_MCA_Review1_Final_v3.pptx')
print("Successfully generated PresenceIQ_MCA_Review1_Final_v3.pptx (15 slides) - Removed specified slides and added Device Binding diagram.")
